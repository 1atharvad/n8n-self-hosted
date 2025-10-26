import uuid
from copy import deepcopy
from pathlib import Path

from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parent.parent.parent
TEMPLATE_PATH = Path(BASE_DIR, "n8n_files", "job_news_template.pptx")
FILES_DIR = Path(BASE_DIR, "n8n_files", "ppt_files")
AUDIO_DIR = Path(BASE_DIR, "n8n_files", "audio_files")


class PPTGenerator:
    def __init__(self):
        """
        A utility class for generating PowerPoint presentations from a
        predefined template. It supports copying, inserting, and removing
        slides while dynamically replacing text placeholders with job-specific
        content.
        """
        self.files_dir = FILES_DIR
        self.files_dir.mkdir(exist_ok=True)
        self.template_slide = 0

    @staticmethod
    def copy_slide(prs, source_slide):
        """
        Duplicates a slide from the source presentation and inserts it into the
        target presentation.

        Args:
            prs (Presentation): The target PowerPoint presentation object.
            source_slide (Slide): The slide to be copied.

        Returns:
            A new slide object that is a duplicate of the source slide.
        """
        new_slide = prs.slides.add_slide(source_slide.slide_layout)
        for shape in source_slide.shapes:
            el = deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
        return new_slide

    def insert_slide(self, prs, template_slide, job, index: int):
        """
        Insert a new slide into the presentation by copying a template slide
        and replacing placeholder text with job-specific content.

        Args:
            prs (Presentation): The target PowerPoint presentation object.
            template_slide (Slide): The slide used as a template.
            job (dict): A dictionary containing job-specific data for text
                replacement.
            index (int): The position in the slide deck where the new slide
                will be inserted.
        """
        new_slide = self.copy_slide(prs, template_slide)
        replacements = {
            value: job[key] for key, value in self.old_text.items()
        }

        for shape in new_slide.shapes:
            if shape.has_text_frame:
                for old_text, new_text in replacements.items():
                    if old_text in shape.text:
                        shape.text = shape.text.replace(old_text, new_text)

        sld_id_lst = prs.slides._sldIdLst
        new_id = sld_id_lst[-1]
        sld_id_lst.remove(new_id)
        sld_id_lst.insert(index, new_id)

    @staticmethod
    def remove_slide(prs, slide):
        """
        Removes a slide from the presentation.

        Args:
            prs (Presentation): The target PowerPoint presentation object.
            slide (Slide): The slide to be removed.
        """
        slide_id_lst = prs.slides._sldIdLst
        slide_index = prs.slides.index(slide)
        slide_id_lst.remove(slide_id_lst[slide_index])

    def create_slide(self, jobs: list[dict]) -> Path:
        """
        Generates a PowerPoint presentation by duplicating a template slide for
        each job entry and replacing placeholders with the corresponding data.

        Args:
            jobs (list[dict]): A list of dictionaries where each dictionary
                represents job-specific data to populate the slides.

        Returns:
            The file path of the generated PowerPoint presentation.
        """
        prs = Presentation(TEMPLATE_PATH)
        template_slide = prs.slides[self.template_slide]

        for index, job in enumerate(jobs):
            self.insert_slide(
                prs, template_slide, job, index + 1 + self.template_slide
            )

        self.remove_slide(prs, template_slide)

        job_id = str(uuid.uuid4())
        file_path = Path(self.files_dir, f"{job_id}.pptx")
        prs.save(file_path)
        return file_path


if __name__ == '__main__':
    ppt = PPTGenerator()
    ppt.template_slide = 0
    ppt.old_text = {
        'company_name': "{{COMPANY_NAME}}",
        'position': "{{POSITION}}",
        'location': "{{LOCATION}}",
        'experience': "{{EXPERIENCE}}",
        'skills': "{{SKILLS}}",
        'job_type': "{{JOB_TYPE}}",
        'link': "{{LINK}}",
    }
    ppt.create_slide(
        [
            {
                "company_name": "Avesta Computer Services",
                "position": "Junior Peoplesoft Developer - Remote",
                "location": "United States",
                'experience': "Entry-level",
                'skills': "PeopleSoft development, Oracle PeopleSoft",
                'job_type': "Full-time, Remote",
                "link": "https://www.linkedin.com/jobs/view/junior-peoplesoft-developer-remote-at-avesta-computer-services-4285538881",
                "audio_file": "9278e699-02d4-466c-a196-96851711a8be.wav",
            },
            {
                "company_name": "Hireshire",
                "position": "Machine Learning Intern",
                "location": "United States",
                'experience': 'None (internship)',
                'skills': 'Basic knowledge of machine learning and programming',
                'job_type': 'Internship',
                "link": "https://www.linkedin.com/jobs/view/machine-learning-intern-at-hireshire-4285519746",
                "audio_file": "9300b75b-cd61-4bdc-bec2-08892c8db469.wav",
            },
        ]
    )
